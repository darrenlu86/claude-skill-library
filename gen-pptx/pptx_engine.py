# -*- coding: utf-8 -*-
"""
PPTX Engine v4 — Generic reusable presentation generator.

Font rules:
  - 28pt: slide title (placeholder)
  - 20pt: section headers within body
  - 16pt: ALL body text (MINIMUM for readable content)
  - 12pt: ONLY sequence diagram arrow labels (exception)

Layout rules:
  - CONTENT_TOP=1.20", CONTENT_BOTTOM=6.65", CONTENT_H=5.45"
  - All builders auto-fill or vertically-center within the safe zone

Configuration:
  - template: pass to PptxEngine(template=...) or None for blank
  - footer_text: pass to PptxEngine(footer_text=...) or None to disable
"""
import os
from datetime import date as _date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════════════════════════════════
BLUE      = RGBColor(0x25, 0x63, 0xEB)
BLUE_LT   = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_MID  = RGBColor(0x93, 0xBB, 0xFD)
ORANGE    = RGBColor(0xF9, 0x73, 0x16)
ORANGE_LT = RGBColor(0xFF, 0xED, 0xD5)
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
GREEN_LT  = RGBColor(0xDC, 0xFC, 0xE7)
RED       = RGBColor(0xEF, 0x44, 0x44)
RED_LT    = RGBColor(0xFE, 0xE2, 0xE2)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LT = RGBColor(0xED, 0xE9, 0xFE)
AMBER     = RGBColor(0xD9, 0x77, 0x06)
AMBER_LT  = RGBColor(0xFE, 0xF3, 0xC7)
CYAN      = RGBColor(0x06, 0xB6, 0xD4)
CYAN_LT   = RGBColor(0xCC, 0xFB, 0xF1)
DARK      = RGBColor(0x1E, 0x29, 0x3B)
GRAY      = RGBColor(0x64, 0x74, 0x8B)
GRAY_LT   = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_BG   = RGBColor(0xF8, 0xFA, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x33, 0x33, 0x33)

FONT = "Microsoft JhengHei"

# ═══════════════════════════════════════════════════════════════
# LAYOUT CONSTANTS (16:9 slide = 13.33" × 7.50")
# ═══════════════════════════════════════════════════════════════
SLIDE_W = 13.33
SLIDE_H = 7.50

MARGIN_X = 0.55
CONTENT_W = 12.20

CONTENT_TOP = 1.20
CONTENT_BOTTOM = 6.65
CONTENT_H = CONTENT_BOTTOM - CONTENT_TOP  # 5.45"

CONTENT_MID_Y = CONTENT_TOP + CONTENT_H / 2

COL2_W = 5.85
COL2_GAP = 0.50
COL2_LEFT_X = MARGIN_X
COL2_RIGHT_X = MARGIN_X + COL2_W + COL2_GAP

# ═══════════════════════════════════════════════════════════════
# FONT SIZE HIERARCHY (STRICT)
#
#   28pt  — slide title (placeholder, set by engine)
#   20pt  — section header inside slide body
#   16pt  — ALL body text (MINIMUM for readable content)
#   12pt  — sequence diagram arrow labels ONLY
#
# ═══════════════════════════════════════════════════════════════
SZ_TITLE   = Pt(28)
SZ_SECTION = Pt(20)
SZ_BODY    = Pt(16)   # ← minimum for all readable text
SZ_SEQ     = Pt(12)   # ← ONLY for sequence diagram labels

# Legacy aliases — all map to SZ_BODY to prevent accidental small text
SZ_SUB   = SZ_BODY
SZ_SMALL = SZ_BODY
SZ_TINY  = SZ_BODY


# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def bh(n, sz=16):
    """Box height (inches) for n lines at sz pt. CJK 16pt ≈ 0.30"/line."""
    line_h = sz * 1.4 / 72
    return round(n * line_h + 0.20, 2)


import unicodedata

def _char_width(ch):
    """Relative width of a character. CJK ≈ 1.0, Latin ≈ 0.55."""
    cat = unicodedata.east_asian_width(ch)
    return 1.0 if cat in ('W', 'F') else 0.55


def _text_width_pt(text, sz_pt):
    """Estimated rendered width of a single line in points."""
    return sum(_char_width(c) for c in text) * sz_pt


def measure_text(text, sz_pt=16, box_w_in=None, padding_in=0.20):
    """Measure text and return (line_count, height_inches).

    If box_w_in is given, estimates how many lines the text wraps into.
    Handles explicit \\n newlines as forced breaks.

    Returns:
        (n_lines, height_in)
    """
    lines = text.split('\n')
    total_lines = 0

    if box_w_in is not None:
        usable_w_pt = (box_w_in - padding_in) * 72  # inches → points
        for line in lines:
            if not line.strip():
                total_lines += 0.6  # blank line counts as partial
                continue
            w = _text_width_pt(line, sz_pt)
            n = max(1, -(-int(w) // int(usable_w_pt)))  # ceil division
            total_lines += n
    else:
        total_lines = len(lines)

    line_h = sz_pt * 1.4 / 72
    height = total_lines * line_h + padding_in
    return (int(total_lines + 0.5), round(height, 2))


_COLOR_MAP = {
    "BLUE": BLUE, "BLUE_LT": BLUE_LT, "BLUE_MID": BLUE_MID,
    "GREEN": GREEN, "GREEN_LT": GREEN_LT,
    "ORANGE": ORANGE, "ORANGE_LT": ORANGE_LT,
    "RED": RED, "RED_LT": RED_LT,
    "PURPLE": PURPLE, "PURPLE_LT": PURPLE_LT,
    "AMBER": AMBER, "AMBER_LT": AMBER_LT,
    "CYAN": CYAN, "CYAN_LT": CYAN_LT,
    "DARK": DARK, "GRAY": GRAY, "GRAY_LT": GRAY_LT, "GRAY_BG": GRAY_BG,
    "WHITE": WHITE, "BLACK": BLACK,
}


def resolve_named_color(val):
    """Resolve a color string name to RGBColor (used internally by builders)."""
    if val is None:
        return BLUE
    if isinstance(val, str):
        upper = val.upper().replace(" ", "_")
        if upper in _COLOR_MAP:
            return _COLOR_MAP[upper]
        h = val.lstrip("#")
        if len(h) == 6:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return val


def _is_dark(color):
    if not isinstance(color, RGBColor):
        return True
    r = int(str(color)[0:2], 16)
    g = int(str(color)[2:4], 16)
    b = int(str(color)[4:6], 16)
    return (0.299 * r + 0.587 * g + 0.114 * b) < 140


def _auto_tc(bg):
    return WHITE if _is_dark(bg) else BLACK


# ═══════════════════════════════════════════════════════════════
# DECORATION MOTIFS
#
# Each motif is a list of (shape_type, dx, dy, w, h) relative offsets.
# dx/dy are relative to the anchor corner. All units in inches.
# Colors are applied at render time (using _LT palette variants).
# ═══════════════════════════════════════════════════════════════

_MOTIF_DOTS = [
    # Scattered circles (abstract data/nodes)
    (MSO_SHAPE.OVAL, 0.00, 0.00, 0.30, 0.30),
    (MSO_SHAPE.OVAL, 0.40, 0.25, 0.20, 0.20),
    (MSO_SHAPE.OVAL, 0.15, 0.55, 0.25, 0.25),
    (MSO_SHAPE.OVAL, 0.50, 0.60, 0.15, 0.15),
]

_MOTIF_BARS = [
    # Abstract bar chart
    (MSO_SHAPE.RECTANGLE, 0.00, 0.50, 0.18, 0.40),
    (MSO_SHAPE.RECTANGLE, 0.22, 0.30, 0.18, 0.60),
    (MSO_SHAPE.RECTANGLE, 0.44, 0.10, 0.18, 0.80),
    (MSO_SHAPE.RECTANGLE, 0.66, 0.35, 0.18, 0.55),
]

_MOTIF_RINGS = [
    # Concentric arcs (abstract process/cycle)
    (MSO_SHAPE.OVAL, 0.00, 0.00, 0.80, 0.80),
    (MSO_SHAPE.OVAL, 0.15, 0.15, 0.50, 0.50),
    (MSO_SHAPE.OVAL, 0.30, 0.30, 0.20, 0.20),
]

_MOTIF_GRID = [
    # Small grid (abstract matrix/table)
    (MSO_SHAPE.ROUNDED_RECTANGLE, 0.00, 0.00, 0.30, 0.30),
    (MSO_SHAPE.ROUNDED_RECTANGLE, 0.35, 0.00, 0.30, 0.30),
    (MSO_SHAPE.ROUNDED_RECTANGLE, 0.00, 0.35, 0.30, 0.30),
    (MSO_SHAPE.ROUNDED_RECTANGLE, 0.35, 0.35, 0.30, 0.30),
]

_MOTIF_ARROW = [
    # Abstract forward arrow (process/flow)
    (MSO_SHAPE.RECTANGLE, 0.00, 0.25, 0.60, 0.20),
    (MSO_SHAPE.ISOSCELES_TRIANGLE, 0.50, 0.10, 0.35, 0.50),
]

# Map slide types → appropriate motif + light color
_SLIDE_MOTIF_MAP = {
    "bullet":     (_MOTIF_DOTS,  BLUE_LT),
    "arch":       (_MOTIF_BARS,  BLUE_LT),
    "flow":       (_MOTIF_ARROW, GREEN_LT),
    "dual_flow":  (_MOTIF_ARROW, GREEN_LT),
    "steps":      (_MOTIF_DOTS,  BLUE_LT),
    "cards":      (_MOTIF_BARS,  CYAN_LT),
    "stat":       (_MOTIF_BARS,  CYAN_LT),
    "table":      (_MOTIF_GRID,  GRAY_LT),
    "highlight":  (_MOTIF_RINGS, PURPLE_LT),
    "timeline":   (_MOTIF_ARROW, ORANGE_LT),
    "comparison": (_MOTIF_DOTS,  GREEN_LT),
    "quote":      (_MOTIF_RINGS, BLUE_LT),
    "matrix":     (_MOTIF_GRID,  AMBER_LT),
    "tree":       (_MOTIF_GRID,  GRAY_LT),
    "sequence":   (_MOTIF_DOTS,  PURPLE_LT),
}


# ═══════════════════════════════════════════════════════════════
# ENGINE
# ═══════════════════════════════════════════════════════════════

def _load_config():
    """Load user config from config.json next to this file.
    Returns dict with keys: template, footer_text, output_dir.
    Missing keys default to None.
    """
    import json as _json
    cfg_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "config.json")
    if os.path.exists(cfg_path):
        with open(cfg_path, "r", encoding="utf-8") as f:
            return _json.load(f)
    return {}

_USER_CONFIG = _load_config()


class PptxEngine:

    def __init__(self, template=None, footer_text=None):
        """Initialize the engine.

        Priority for template/footer:
          1. Explicit argument (from JSON "template"/"footer_text")
          2. User config (~/.claude/skills/gen-pptx/config.json)
          3. Blank presentation (no template, no footer)

        Args:
            template: Path to .pptx template.
                      None = use config.json default.
                      "" = force blank (no template).
            footer_text: Copyright footer text.
                         None = use config.json default.
                         "" = force no footer.
        """
        # Resolve template
        resolved_template = template
        if resolved_template is None:
            resolved_template = _USER_CONFIG.get("template")

        if resolved_template == "" or resolved_template is None:
            self.prs = Presentation()
        else:
            self.prs = Presentation(resolved_template)

        # Resolve footer
        self._footer_text = footer_text
        if self._footer_text is None:
            self._footer_text = _USER_CONFIG.get("footer_text")
        self._remove_existing_slides()
        n_layouts = len(self.prs.slide_layouts)
        self.L_TITLE   = self.prs.slide_layouts[0] if n_layouts > 0 else None
        self.L_SECTION = self.prs.slide_layouts[min(1, n_layouts - 1)]
        self.L_TONLY   = self.prs.slide_layouts[min(2, n_layouts - 1)]
        self.L_CONTENT = self.prs.slide_layouts[min(3, n_layouts - 1)]
        self.L_PLAIN   = self.prs.slide_layouts[min(6, n_layouts - 1)]
        self.L_BLANK   = self.prs.slide_layouts[min(9, n_layouts - 1)]
        self._page = 0
        self.auto_decorate = True  # enable/disable decoration

    def _remove_existing_slides(self):
        for sldId in list(self.prs.slides._sldIdLst):
            rId = sldId.get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId:
                self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(sldId)

    def _next_page(self):
        self._page += 1
        return self._page

    def _new_slide(self, title, layout=None):
        slide = self.prs.slides.add_slide(layout or self.L_PLAIN)
        page = self._next_page()
        self._set_title(slide, title)
        self._set_footer(slide, page)
        return slide

    def _set_title(self, slide, text, size=SZ_TITLE, color=DARK):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                shape.text = ""
                p = shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.size = size
                run.font.bold = True
                run.font.color.rgb = color
                run.font.name = FONT
                return

    def _set_footer(self, slide, page=None):
        if page is None:
            page = self._page
        today_str = _date.today().strftime("%Y/%m/%d")

        # Footer placeholders (idx 10/11/12) are often NOT inherited from
        # layout. Clone them from the slide layout if missing; if the
        # layout itself lacks them, fall back to any layout that has them.
        existing_idx = {ph.placeholder_format.idx for ph in slide.placeholders}
        needed = {10, 11, 12} - existing_idx
        if needed:
            from copy import deepcopy
            # Try current layout first, then fall back to others
            source_layout = slide.slide_layout
            source_phs = {ph.placeholder_format.idx: ph for ph in source_layout.placeholders}
            if not (needed & set(source_phs.keys())):
                for layout in self.prs.slide_layouts:
                    candidate = {ph.placeholder_format.idx: ph for ph in layout.placeholders}
                    if needed & set(candidate.keys()):
                        source_phs = candidate
                        break
            for idx_needed in needed:
                if idx_needed in source_phs:
                    sp_clone = deepcopy(source_phs[idx_needed]._element)
                    slide.shapes._spTree.append(sp_clone)

        # Now set values on all footer placeholders
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 10:
                shape.text = today_str
            elif idx == 11:
                shape.text = self._footer_text or ""
            elif idx == 12:
                shape.text = str(page)

    # ═══════════════════════════════════════════════════════════
    # PRIMITIVES (all default to SZ_BODY = 16pt)
    # ═══════════════════════════════════════════════════════════

    def box(self, slide, x, y, w, h, fill, border=None, text="",
            sz=SZ_BODY, bold=False, tc=None, align=PP_ALIGN.CENTER,
            anchor="ctr"):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if border:
            shape.line.color.rgb = border
            shape.line.width = Pt(0.75)
        else:
            shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        bodyPr = tf._txBody.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        pad_v = int(0.06 * 914400)
        pad_h = int(0.10 * 914400)
        bodyPr.set('tIns', str(pad_v))
        bodyPr.set('bIns', str(pad_v))
        bodyPr.set('lIns', str(pad_h))
        bodyPr.set('rIns', str(pad_h))
        bodyPr.set('anchor', anchor)

        if not text:
            return shape

        if tc is None:
            tc = _auto_tc(fill)

        lines = text.split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(2)
            p.space_before = Pt(0)
            if len(lines) > 3:
                p.line_spacing = 1.15
            run = p.add_run()
            run.text = line
            run.font.size = sz
            run.font.bold = bold
            run.font.color.rgb = tc
            run.font.name = FONT

        return shape

    def rect(self, slide, x, y, w, h, fill, border=None, text="",
             sz=SZ_BODY, bold=False, tc=None, align=PP_ALIGN.CENTER):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if border:
            shape.line.color.rgb = border
            shape.line.width = Pt(0.75)
        else:
            shape.line.fill.background()

        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            if tc is None:
                tc = _auto_tc(fill)
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = sz
            run.font.bold = bold
            run.font.color.rgb = tc
            run.font.name = FONT
        return shape

    def textbox(self, slide, x, y, w, h, text, sz=SZ_BODY,
                bold=False, color=BLACK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = sz
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
        return txBox

    def bullet_box(self, slide, x, y, w, h, items, sz=SZ_BODY,
                   color=BLACK, title=None, title_sz=SZ_SECTION,
                   title_color=DARK, line_sp=1.3):
        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True

        first = True
        if title:
            p = tf.paragraphs[0]
            p.line_spacing = line_sp
            run = p.add_run()
            run.text = title
            run.font.size = title_sz
            run.font.bold = True
            run.font.color.rgb = title_color
            run.font.name = FONT
            p.space_after = Pt(8)
            first = False

        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(4)
            p.line_spacing = line_sp

            if isinstance(item, tuple):
                text, level = item
            else:
                text, level = item, 0

            p.level = level
            bullet = "●  " if level == 0 else "    ○  "

            run = p.add_run()
            run.text = bullet + text
            # Sub-items: 16pt → 16pt (no shrink below SZ_BODY)
            run.font.size = sz
            run.font.color.rgb = color
            run.font.name = FONT

        return txBox

    def arrow_text(self, slide, x, y, char="→", color=BLUE, sz=Pt(20)):
        self.textbox(slide, x, y, 0.35, 0.35, char,
                     sz=sz, bold=True, color=color, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════
    # DECORATION — auto-fill sparse slides with geometric motifs
    # ═══════════════════════════════════════════════════════════

    def _content_coverage(self, slide):
        """Estimate what fraction of the content zone is covered by shapes."""
        zone_area = CONTENT_W * CONTENT_H
        covered = 0
        for shape in slide.shapes:
            sx = shape.left / 914400   # EMU → inches
            sy = shape.top / 914400
            sw = shape.width / 914400
            sh = shape.height / 914400
            # Only count shapes within the content zone
            if sy >= CONTENT_TOP - 0.10 and sy + sh <= CONTENT_BOTTOM + 0.30:
                covered += sw * sh
        return min(covered / zone_area, 1.0) if zone_area > 0 else 1.0

    def _find_empty_corner(self, slide):
        """Find the emptiest corner of the content zone.
        Returns (x, y) anchor point for decoration placement.
        Checks 4 quadrants: top-right, bottom-right, bottom-left, top-left.
        """
        corners = [
            ("tr", MARGIN_X + CONTENT_W - 1.2, CONTENT_TOP + 0.10),
            ("br", MARGIN_X + CONTENT_W - 1.2, CONTENT_BOTTOM - 1.1),
            ("bl", MARGIN_X + 0.10,             CONTENT_BOTTOM - 1.1),
            ("tl", MARGIN_X + 0.10,             CONTENT_TOP + 0.10),
        ]
        best = corners[0]  # default: top-right
        min_overlap = float('inf')

        for (name, cx, cy) in corners:
            overlap = 0
            for shape in slide.shapes:
                sx = shape.left / 914400
                sy = shape.top / 914400
                sw = shape.width / 914400
                sh = shape.height / 914400
                # Check overlap with 1.2" × 1.0" corner zone
                ox = max(0, min(cx + 1.2, sx + sw) - max(cx, sx))
                oy = max(0, min(cy + 1.0, sy + sh) - max(cy, sy))
                overlap += ox * oy
            if overlap < min_overlap:
                min_overlap = overlap
                best = (name, cx, cy)

        return best[1], best[2]

    def _place_motif(self, slide, motif, color, anchor_x, anchor_y):
        """Place a decoration motif at the given anchor position."""
        from lxml import etree
        for (shape_type, dx, dy, w, h) in motif:
            shape = slide.shapes.add_shape(
                shape_type,
                Inches(anchor_x + dx), Inches(anchor_y + dy),
                Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            # Set 40% transparency via XML (python-pptx doesn't have native API)
            solidFill = shape.fill._fill
            srgb = solidFill.find(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            if srgb is not None:
                color_elem = srgb[0] if len(srgb) > 0 else None
                if color_elem is not None:
                    color_elem.set('alpha', '40000')  # 40% opacity (40000/100000)

    def decorate(self, slide, slide_type, accent_color=None):
        """Add decorative motif to a slide if it has sparse content.
        Called automatically by builders when self.auto_decorate is True.
        """
        if not self.auto_decorate:
            return
        coverage = self._content_coverage(slide)
        if coverage > 0.45:
            return  # slide is full enough

        motif_info = _SLIDE_MOTIF_MAP.get(slide_type)
        if not motif_info:
            return
        motif, default_color = motif_info
        color = accent_color or default_color
        ax, ay = self._find_empty_corner(slide)
        self._place_motif(slide, motif, color, ax, ay)

    # ═══════════════════════════════════════════════════════════
    # SLIDE BUILDERS
    # ═══════════════════════════════════════════════════════════

    def _title_like_slide(self, layout, title, subtitle="",
                          title_color=DARK, sub_color=GRAY, show_footer=True):
        """Shared builder for title / section / end slides."""
        slide = self.prs.slides.add_slide(layout)
        page = self._next_page()
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 0:
                shape.text = ""
                p = shape.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = title
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = title_color
                run.font.name = FONT
            elif idx == 1 and subtitle:
                shape.text = ""
                p = shape.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = subtitle
                run.font.size = Pt(20)
                run.font.color.rgb = sub_color
                run.font.name = FONT
        if show_footer:
            self._set_footer(slide, page)
        return slide

    def title_slide(self, title, subtitle="", title_color=DARK, sub_color=GRAY):
        return self._title_like_slide(self.L_TITLE, title, subtitle,
                                      title_color=title_color, sub_color=sub_color,
                                      show_footer=False)

    def section_slide(self, title, subtitle="", title_color=DARK, sub_color=GRAY):
        return self._title_like_slide(self.L_SECTION, title, subtitle,
                                      title_color=title_color, sub_color=sub_color)

    def end_slide(self, title="謝謝各位", subtitle="Q&A 時間",
                  title_color=DARK, sub_color=GRAY):
        return self._title_like_slide(self.L_SECTION, title, subtitle,
                                      title_color=title_color, sub_color=sub_color,
                                      show_footer=False)

    # ─── Agenda ─────────────────────────────────────────────

    def agenda_slide(self, title, items):
        slide = self._new_slide(title)
        n = len(items)
        row_h = 0.70
        gap = 0.15
        total_h = n * row_h + (n - 1) * gap
        start_y = CONTENT_TOP + (CONTENT_H - total_h) / 2

        for i, (num, topic, desc) in enumerate(items):
            y = start_y + i * (row_h + gap)
            self.box(slide, 1.0, y, 0.55, 0.55, BLUE, None, str(num),
                     sz=SZ_BODY, bold=True, tc=WHITE)
            self.textbox(slide, 1.8, y - 0.02, 5, 0.38, topic,
                         sz=SZ_SECTION, bold=True, color=DARK)
            self.textbox(slide, 1.8, y + 0.35, 10, 0.32, desc,
                         sz=SZ_BODY, color=GRAY)
        return slide

    # ─── Bullet Slide ───────────────────────────────────────

    def _estimate_bullet_height(self, col, line_h=0.38):
        """Estimate the height needed for a bullet column."""
        h = 0
        if col.get("header"):
            h += 0.50 + 0.12  # header bar + gap
        if col.get("title"):
            h += 0.45  # title line + space
        items = col.get("items", [])
        for item in items:
            if isinstance(item, (tuple, list)):
                h += line_h  # indented item
            else:
                # Rough estimate: long text may wrap
                text = str(item)
                lines = max(1, len(text) // 30)  # ~30 chars per line at 16pt
                h += line_h * lines
        return h

    def bullet_slide(self, title, columns, subtitle=None):
        slide = self._new_slide(title)
        zone_top = CONTENT_TOP
        zone_h = CONTENT_H

        if subtitle:
            self.textbox(slide, MARGIN_X, zone_top, CONTENT_W, 0.38, subtitle,
                         sz=SZ_BODY, color=GRAY)
            zone_top += 0.48
            zone_h -= 0.48

        n = len(columns)
        if n == 1:
            positions = [(MARGIN_X, CONTENT_W)]
        elif n == 2:
            positions = [(COL2_LEFT_X, COL2_W), (COL2_RIGHT_X, COL2_W)]
        else:
            gap = 0.35
            cw = (CONTENT_W - gap * (n - 1)) / n
            positions = [(MARGIN_X + i * (cw + gap), cw) for i in range(n)]

        # Estimate max content height across all columns
        max_h = max(self._estimate_bullet_height(col) for col in columns)
        # Vertically center: offset from zone_top
        y_offset = max(0, (zone_h - max_h) / 2)
        y_start = zone_top + y_offset

        for i, col in enumerate(columns):
            x, w = positions[i]
            cy = y_start

            header = col.get("header")
            if header:
                hc = col.get("header_color", BLUE)
                hh = 0.50
                self.box(slide, x, cy, w, hh, hc, None, header,
                         sz=SZ_BODY, bold=True, tc=WHITE)
                cy += hh + 0.12

            items = col.get("items", [])
            col_title = col.get("title")
            tc = col.get("title_color", DARK)
            bullet_h = CONTENT_BOTTOM - cy
            self.bullet_box(slide, x, cy, w, bullet_h, items,
                            sz=SZ_BODY, title=col_title, title_color=tc)
        return slide

    # ─── Architecture Diagram ───────────────────────────────

    def arch_slide(self, title, layers, note=None):
        slide = self._new_slide(title)
        n = len(layers)
        arrow_w = 0.40
        gap = 0.08
        total_arrows = (n - 1) * (arrow_w + gap * 2)
        box_w = (CONTENT_W - total_arrows) / n

        # Measure note height dynamically
        note_gap = 0.15 if note else 0
        if note:
            _, note_h = measure_text(note, 16, CONTENT_W - 0.40)
            note_h = max(note_h, 0.50)
        else:
            note_h = 0

        # Measure box content to find minimum needed height
        max_text_h = 0
        for layer in layers:
            label = layer[0]
            body = layer[2] if len(layer) > 2 else ""
            full_text = f"{label}\n\n{body}" if body else label
            _, th = measure_text(full_text, 16, box_w - 0.30)
            max_text_h = max(max_text_h, th)

        available_h = CONTENT_H - note_h - note_gap
        box_h = max(max_text_h, min(available_h, 4.5))

        # Vertically center boxes (and note) in the safe zone
        total_content_h = box_h + note_gap + note_h
        box_y = CONTENT_TOP + (CONTENT_H - total_content_h) / 2

        for i, layer in enumerate(layers):
            label = layer[0]
            fill = layer[1]
            body = layer[2] if len(layer) > 2 else ""
            bx = MARGIN_X + i * (box_w + arrow_w + gap * 2)

            full_text = f"{label}\n\n{body}" if body else label
            self.box(slide, bx, box_y, box_w, box_h, fill, None,
                     full_text, sz=SZ_BODY, tc=_auto_tc(fill),
                     align=PP_ALIGN.CENTER)

            if i < n - 1:
                ax = bx + box_w + gap
                ay = box_y + box_h / 2 - 0.18
                self.arrow_text(slide, ax, ay, "→", BLUE, Pt(22))

        if note:
            ny = box_y + box_h + note_gap
            self.box(slide, MARGIN_X, ny, CONTENT_W, note_h, GRAY_LT, None,
                     note, sz=SZ_BODY, tc=GRAY)
        return slide

    # ─── Flow Diagram ───────────────────────────────────────

    def flow_slide(self, title, steps, subtitle=None):
        slide = self._new_slide(title)
        sub_offset = 0
        if subtitle:
            self.textbox(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 0.38,
                         subtitle, sz=SZ_BODY, color=GRAY)
            sub_offset = 0.48

        n = len(steps)
        arrow_w = 0.35
        gap = 0.06
        total_arrows = (n - 1) * (arrow_w + gap * 2)
        sw = (CONTENT_W - total_arrows) / n

        # Measure: find tallest step label to set box height
        max_lines = 1
        for step in steps:
            label = step[0]
            lines, _ = measure_text(label, 16, sw)
            max_lines = max(max_lines, lines)
        sh = max(1.2, min(2.2, bh(max_lines, 16) + 0.30))

        zone_top = CONTENT_TOP + sub_offset
        zone_h = CONTENT_H - sub_offset
        sy = zone_top + (zone_h - sh) / 2

        for i, step in enumerate(steps):
            label = step[0]
            color = step[1] if len(step) > 1 else BLUE
            sx = MARGIN_X + i * (sw + arrow_w + gap * 2)
            self.box(slide, sx, sy, sw, sh, color, None, label,
                     sz=SZ_BODY, tc=_auto_tc(color))
            if i < n - 1:
                ax = sx + sw + gap
                self.arrow_text(slide, ax, sy + sh / 2 - 0.18, "→", BLUE, Pt(22))
        return slide

    # ─── Dual Flow ──────────────────────────────────────────

    def dual_flow_slide(self, title, flows):
        slide = self._new_slide(title)
        n_flows = len(flows)
        label_h = 0.40
        flow_gap = 0.25

        # Measure: find max lines across ALL step labels in ALL flows
        max_lines_per_flow = []
        for flow in flows:
            steps = flow["steps"]
            ns = len(steps)
            arrow_w = 0.30
            gap_f = 0.05
            total_arrows = (ns - 1) * (arrow_w + gap_f * 2)
            sw = (CONTENT_W - total_arrows) / ns
            ml = 1
            for (label, _bg) in steps:
                lines, _ = measure_text(label, 16, sw)
                ml = max(ml, lines)
            max_lines_per_flow.append(ml)

        # Compute row heights: enough for content, but capped and distributed
        min_row_h = max(bh(m, 16) + 0.20 for m in max_lines_per_flow)
        total_labels = n_flows * label_h
        total_gaps = (n_flows - 1) * flow_gap + n_flows * 0.10
        max_row_h = (CONTENT_H - total_labels - total_gaps) / n_flows
        row_h = max(min_row_h, min(max_row_h, 2.2))

        total_h = n_flows * (label_h + 0.10 + row_h) + (n_flows - 1) * flow_gap
        y = CONTENT_TOP + (CONTENT_H - total_h) / 2

        for fi, flow in enumerate(flows):
            lc = flow.get("label_color", BLUE)
            self.textbox(slide, MARGIN_X, y, 8, label_h, flow["label"],
                         sz=SZ_SECTION, bold=True, color=lc)
            y += label_h + 0.10

            steps = flow["steps"]
            ns = len(steps)
            arrow_w = 0.30
            gap = 0.05
            total_arrows = (ns - 1) * (arrow_w + gap * 2)
            sw = (CONTENT_W - total_arrows) / ns

            for i, (label, bg) in enumerate(steps):
                sx = MARGIN_X + i * (sw + arrow_w + gap * 2)
                self.box(slide, sx, y, sw, row_h, bg, None, label,
                         sz=SZ_BODY, tc=_auto_tc(bg))
                if i < ns - 1:
                    ax = sx + sw + gap
                    self.arrow_text(slide, ax, y + row_h / 2 - 0.15,
                                    "→", lc, Pt(18))
            y += row_h + flow_gap
        return slide

    # ─── Steps ──────────────────────────────────────────────

    def _circle(self, slide, x, y, size, fill, text="",
                sz=SZ_SECTION, bold=True, tc=WHITE):
        """Draw a circle with centered text."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y), Inches(size), Inches(size))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = False
        bodyPr = tf._txBody.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        bodyPr.set('anchor', 'ctr')
        pad = int(0.02 * 914400)
        bodyPr.set('tIns', str(pad))
        bodyPr.set('bIns', str(pad))
        bodyPr.set('lIns', str(pad))
        bodyPr.set('rIns', str(pad))

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = sz
        run.font.bold = bold
        run.font.color.rgb = tc
        run.font.name = FONT
        return shape

    def steps_slide(self, title, steps, notes=None, notes_title="注意事項"):
        slide = self._new_slide(title)
        has_notes = notes is not None
        main_w = 7.2 if has_notes else CONTENT_W

        n = len(steps)
        badge_sz = 0.50
        text_w = main_w - badge_sz - 0.20
        arrow_h = 0.30

        # Measure each step to get dynamic content height
        step_heights = []
        for (step_title, step_desc) in steps:
            _, title_h = measure_text(step_title, 20, text_w)
            _, desc_h = measure_text(step_desc, 16, text_w)
            sh = title_h + desc_h
            step_heights.append(max(sh, badge_sz + 0.10))

        total_h = sum(step_heights) + (n - 1) * arrow_h
        # If it doesn't fit, scale down proportionally
        if total_h > CONTENT_H:
            scale = CONTENT_H / total_h
            step_heights = [h * scale for h in step_heights]
            arrow_h *= scale
            total_h = CONTENT_H

        start_y = CONTENT_TOP + (CONTENT_H - total_h) / 2

        y = start_y
        for i, (step_title, step_desc) in enumerate(steps):
            sh = step_heights[i]
            _, title_h = measure_text(step_title, 20, text_w)
            desc_h = sh - title_h

            # Numbered circle badge — vertically centered to step
            badge_y = y + (sh - badge_sz) / 2
            self._circle(slide, MARGIN_X, badge_y, badge_sz, BLUE,
                         str(i + 1), sz=SZ_SECTION, tc=WHITE)
            # Title and description
            text_x = MARGIN_X + badge_sz + 0.20
            self.textbox(slide, text_x, y,
                         text_w, title_h,
                         step_title, sz=SZ_SECTION, bold=True, color=DARK)
            self.textbox(slide, text_x, y + title_h,
                         text_w, desc_h,
                         step_desc, sz=SZ_BODY, color=GRAY)

            if i < n - 1:
                ay = y + sh
                arrow_x = MARGIN_X + badge_sz / 2 - 0.10
                self.textbox(slide, arrow_x, ay, 0.30, arrow_h, "↓",
                             sz=Pt(22), bold=True, color=BLUE,
                             align=PP_ALIGN.CENTER)
            y += sh + arrow_h

        if has_notes:
            items = notes if isinstance(notes, list) else notes.get("items", [])
            n_title = notes_title if isinstance(notes, list) else notes.get("title", notes_title)
            panel_x = 8.3
            panel_w = CONTENT_W - panel_x + MARGIN_X
            panel_h = CONTENT_H
            self.box(slide, panel_x, CONTENT_TOP, panel_w, panel_h,
                     AMBER_LT, AMBER, "", anchor="t")
            self.bullet_box(slide, panel_x + 0.15, CONTENT_TOP + 0.15,
                            panel_w - 0.30, panel_h - 0.30,
                            items, sz=SZ_BODY, title=n_title,
                            title_sz=SZ_SECTION, title_color=AMBER)
        return slide

    # ─── Tree ───────────────────────────────────────────────

    def tree_slide(self, title, tree_text, font_size=16):
        slide = self._new_slide(title)
        # Enforce minimum 16pt (except allow 14pt for dense trees)
        font_size = max(font_size, 14)
        n_lines = tree_text.count('\n') + 1
        text_h = bh(n_lines, font_size)
        text_h = min(text_h, CONTENT_H)
        y = CONTENT_TOP + (CONTENT_H - text_h) / 2

        self.box(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, CONTENT_H,
                 GRAY_BG, GRAY_LT, "", anchor="t")
        self.textbox(slide, MARGIN_X + 0.25, y, CONTENT_W - 0.50, text_h,
                     tree_text, sz=Pt(font_size), color=BLACK)
        return slide

    # ─── Cards ──────────────────────────────────────────────

    def cards_slide(self, title, cards, subtitle=None):
        slide = self._new_slide(title)
        sub_offset = 0
        if subtitle:
            self.textbox(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 0.38,
                         subtitle, sz=SZ_BODY, color=GRAY)
            sub_offset = 0.50

        n = len(cards)
        gap = 0.25
        cw = (CONTENT_W - gap * (n - 1)) / n
        ch = min(2.0, (CONTENT_H - sub_offset) * 0.50)

        zone_top = CONTENT_TOP + sub_offset
        zone_h = CONTENT_H - sub_offset
        cy = zone_top + (zone_h - ch) / 2

        for i, card in enumerate(cards):
            label = card[0]
            value = card[1]
            accent = card[2] if len(card) > 2 else BLUE
            cx = MARGIN_X + i * (cw + gap)
            self.box(slide, cx, cy, cw, ch, WHITE, GRAY_LT,
                     f"{label}\n\n{value}", sz=SZ_BODY, tc=DARK,
                     align=PP_ALIGN.CENTER)
            self.rect(slide, cx, cy, cw, 0.06, accent)
        return slide

    # ─── Table ──────────────────────────────────────────────

    def table_slide(self, title, headers, rows, col_widths=None,
                    header_color=BLUE, stripe=True):
        slide = self._new_slide(title)
        n_cols = len(headers)
        n_rows = len(rows) + 1
        # Measure longest cell text to set row height
        cw_avg = CONTENT_W / n_cols
        max_lines = 1
        for row in rows:
            for val in row:
                lines, _ = measure_text(str(val), 16, cw_avg - 0.20)
                max_lines = max(max_lines, lines)
        min_row_h = bh(max_lines, 16)
        row_h = max(min_row_h, min(0.50, CONTENT_H / n_rows))
        table_h = min(n_rows * row_h, CONTENT_H)
        table_w = CONTENT_W
        ty = CONTENT_TOP + (CONTENT_H - table_h) / 2

        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(MARGIN_X), Inches(ty),
            Inches(table_w), Inches(table_h))
        tbl = tbl_shape.table

        if col_widths:
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Inches(w)
        else:
            cw = table_w / n_cols
            for i in range(n_cols):
                tbl.columns[i].width = Inches(cw)

        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            self._style_cell(cell, header_color, WHITE, bold=True)

        for i, row in enumerate(rows):
            bg = GRAY_BG if stripe and i % 2 == 1 else WHITE
            for j, val in enumerate(row):
                cell = tbl.cell(i + 1, j)
                cell.text = str(val)
                self._style_cell(cell, bg, BLACK)
        return slide

    def _style_cell(self, cell, fill_color, text_color, bold=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = SZ_BODY
                run.font.color.rgb = text_color
                run.font.bold = bold
                run.font.name = FONT

    # ─── Sequence Diagram ───────────────────────────────────

    def sequence_slide(self, title, actors, messages):
        slide = self._new_slide(title)
        n = len(actors)
        gap = 0.30
        aw = (CONTENT_W - gap * (n - 1)) / n
        ah = 0.55
        actor_y = CONTENT_TOP

        actor_cx = []
        for i, (name, color) in enumerate(actors):
            ax = MARGIN_X + i * (aw + gap)
            actor_cx.append(ax + aw / 2)
            self.box(slide, ax, actor_y, aw, ah, color, None, name,
                     sz=SZ_BODY, bold=True, tc=_auto_tc(color))

        line_top = actor_y + ah
        line_bottom = CONTENT_BOTTOM - 0.10
        for cx in actor_cx:
            self.rect(slide, cx - 0.01, line_top, 0.02,
                      line_bottom - line_top, GRAY_LT)

        n_msgs = len(messages)
        msg_zone = line_bottom - line_top - 0.20
        msg_gap = msg_zone / max(n_msgs, 1)

        for mi, (fi, ti, label) in enumerate(messages):
            my = line_top + 0.10 + mi * msg_gap
            x1 = actor_cx[fi]
            x2 = actor_cx[ti]
            left_x = min(x1, x2) + 0.05
            right_x = max(x1, x2) - 0.05
            arr_w = right_x - left_x

            self.rect(slide, left_x, my + 0.14, arr_w, 0.02, GRAY)

            arrow_char = "→" if x2 > x1 else "←"
            arrow_x = right_x - 0.15 if x2 > x1 else left_x - 0.05
            self.textbox(slide, arrow_x, my - 0.02, 0.30, 0.30,
                         arrow_char, sz=SZ_BODY, bold=True, color=BLUE,
                         align=PP_ALIGN.CENTER)

            # Sequence labels are the ONLY exception: 12pt
            label_x = left_x + arr_w / 2 - 1.0
            self.textbox(slide, label_x, my - 0.18, 2.0, 0.28,
                         label, sz=SZ_SEQ, color=DARK, align=PP_ALIGN.CENTER)
        return slide

    # ─── Highlight ──────────────────────────────────────────

    def highlight_slide(self, title, points):
        slide = self._new_slide(title)
        n = len(points)
        gap = 0.20

        # Measure each point to determine box height
        point_heights = []
        for (heading, desc) in points:
            text = f"{heading}\n{desc}"
            _, h = measure_text(text, 16, CONTENT_W - 0.40)
            point_heights.append(max(h, 0.70))

        total_h = sum(point_heights) + (n - 1) * gap
        # If total fits, distribute evenly for visual balance
        if total_h <= CONTENT_H:
            # Use even height for cleaner look, but at least measured minimum
            even_h = (CONTENT_H - gap * (n - 1)) / n
            box_heights = [max(ph, min(even_h, 1.8)) for ph in point_heights]
        else:
            # Scale down to fit
            scale = (CONTENT_H - gap * (n - 1)) / sum(point_heights)
            box_heights = [max(ph * scale, 0.60) for ph in point_heights]

        total_h = sum(box_heights) + (n - 1) * gap
        start_y = CONTENT_TOP + (CONTENT_H - total_h) / 2

        y = start_y
        for i, (heading, desc) in enumerate(points):
            self.box(slide, MARGIN_X, y, CONTENT_W, box_heights[i], DARK, None,
                     f"{heading}\n{desc}", sz=SZ_BODY, tc=WHITE,
                     align=PP_ALIGN.LEFT, anchor="ctr")
            y += box_heights[i] + gap
        return slide

    # ─── Timeline ──────────────────────────────────────────

    def timeline_slide(self, title, milestones, subtitle=None):
        """Horizontal timeline with milestones.
        milestones: [(label, description, color), ...]
        """
        slide = self._new_slide(title)
        sub_offset = 0
        if subtitle:
            self.textbox(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 0.38,
                         subtitle, sz=SZ_BODY, color=GRAY)
            sub_offset = 0.50

        n = len(milestones)
        zone_top = CONTENT_TOP + sub_offset
        zone_h = CONTENT_H - sub_offset

        # Line in the middle
        line_y = zone_top + zone_h * 0.40
        self.rect(slide, MARGIN_X + 0.3, line_y, CONTENT_W - 0.6, 0.04, GRAY_LT)

        # Milestones evenly distributed
        gap = CONTENT_W / (n + 1)
        dot_r = 0.22
        card_w = min(2.2, gap - 0.20)
        card_h = 1.6

        for i, ms in enumerate(milestones):
            label = ms[0]
            desc = ms[1] if len(ms) > 1 else ""
            color = ms[2] if len(ms) > 2 else BLUE

            cx = MARGIN_X + gap * (i + 1)

            # Dot on the line
            self.box(slide, cx - dot_r / 2, line_y - dot_r / 2 + 0.02,
                     dot_r, dot_r, color, None, "", anchor="ctr")

            # Card: alternate above/below the line
            if i % 2 == 0:
                cy = line_y - card_h - 0.25
            else:
                cy = line_y + 0.30

            self.box(slide, cx - card_w / 2, cy, card_w, card_h,
                     WHITE, color, f"{label}\n\n{desc}",
                     sz=SZ_BODY, tc=DARK, align=PP_ALIGN.CENTER, anchor="ctr")
        return slide

    # ─── Comparison ────────────────────────────────────────

    def comparison_slide(self, title, left, right, subtitle=None):
        """Side-by-side comparison.
        left/right: {"header": "...", "header_color": ..., "items": [...]}
        """
        slide = self._new_slide(title)
        zone_top = CONTENT_TOP
        if subtitle:
            self.textbox(slide, MARGIN_X, zone_top, CONTENT_W, 0.38,
                         subtitle, sz=SZ_BODY, color=GRAY)
            zone_top += 0.50

        col_w = 5.70
        gap = 0.80
        left_x = MARGIN_X
        right_x = MARGIN_X + col_w + gap
        header_h = 0.55
        zone_h = CONTENT_BOTTOM - zone_top

        # VS divider in the middle
        vs_x = left_x + col_w + (gap - 0.50) / 2
        vs_y = zone_top + zone_h / 2 - 0.25
        self.box(slide, vs_x, vs_y, 0.50, 0.50, GRAY, None, "VS",
                 sz=SZ_BODY, bold=True, tc=WHITE)

        for i, col in enumerate([left, right]):
            x = left_x if i == 0 else right_x
            cy = zone_top
            hc = col.get("header_color", BLUE if i == 0 else GREEN)
            header = col.get("header", "")

            # Header bar
            self.box(slide, x, cy, col_w, header_h, hc, None, header,
                     sz=SZ_BODY, bold=True, tc=WHITE)
            cy += header_h + 0.12

            # Items
            items = col.get("items", [])
            bullet_h = CONTENT_BOTTOM - cy
            self.bullet_box(slide, x + 0.10, cy, col_w - 0.20, bullet_h,
                            items, sz=SZ_BODY)
        return slide

    # ─── Quote ─────────────────────────────────────────────

    def quote_slide(self, title, quote_text, author="", accent_color=BLUE):
        """Centered quote with large text."""
        slide = self._new_slide(title)

        # Quote mark
        self.textbox(slide, MARGIN_X + 0.5, CONTENT_TOP + 0.3, 1.0, 0.8,
                     "\u201C", sz=Pt(72), bold=True, color=accent_color)

        # Quote text centered
        quote_h = bh(quote_text.count('\n') + 2, 20)
        quote_h = min(quote_h, 3.0)
        qy = CONTENT_TOP + (CONTENT_H - quote_h) / 2 - 0.2
        self.textbox(slide, MARGIN_X + 1.2, qy, CONTENT_W - 2.4, quote_h,
                     quote_text, sz=SZ_SECTION, color=DARK,
                     align=PP_ALIGN.CENTER)

        # Author
        if author:
            ay = qy + quote_h + 0.25
            self.textbox(slide, MARGIN_X + 1.2, ay, CONTENT_W - 2.4, 0.40,
                         f"— {author}", sz=SZ_BODY, color=GRAY,
                         align=PP_ALIGN.RIGHT)
        return slide

    # ─── Matrix (2×2 grid) ─────────────────────────────────

    def matrix_slide(self, title, cells, x_labels=None, y_labels=None,
                     subtitle=None):
        """2×2 matrix grid.
        cells: [[top_left, top_right], [bottom_left, bottom_right]]
          Each cell: {"label": "...", "desc": "...", "color": "BLUE_LT"}
        x_labels: [left_label, right_label] (optional axis labels)
        y_labels: [top_label, bottom_label] (optional axis labels)
        """
        slide = self._new_slide(title)
        zone_top = CONTENT_TOP
        if subtitle:
            self.textbox(slide, MARGIN_X, zone_top, CONTENT_W, 0.38,
                         subtitle, sz=SZ_BODY, color=GRAY)
            zone_top += 0.50

        # Reserve space for labels
        label_w = 0.50 if y_labels else 0
        label_h = 0.40 if x_labels else 0

        grid_x = MARGIN_X + label_w
        grid_y = zone_top + label_h
        grid_w = CONTENT_W - label_w
        grid_h = CONTENT_BOTTOM - grid_y
        cell_gap = 0.15
        cell_w = (grid_w - cell_gap) / 2
        cell_h = (grid_h - cell_gap) / 2

        # Axis labels
        if x_labels and len(x_labels) >= 2:
            self.textbox(slide, grid_x, zone_top, cell_w, label_h,
                         x_labels[0], sz=SZ_BODY, bold=True, color=GRAY,
                         align=PP_ALIGN.CENTER)
            self.textbox(slide, grid_x + cell_w + cell_gap, zone_top,
                         cell_w, label_h,
                         x_labels[1], sz=SZ_BODY, bold=True, color=GRAY,
                         align=PP_ALIGN.CENTER)

        if y_labels and len(y_labels) >= 2:
            self.textbox(slide, MARGIN_X, grid_y, label_w, cell_h,
                         y_labels[0], sz=SZ_BODY, bold=True, color=GRAY,
                         align=PP_ALIGN.CENTER)
            self.textbox(slide, MARGIN_X, grid_y + cell_h + cell_gap,
                         label_w, cell_h,
                         y_labels[1], sz=SZ_BODY, bold=True, color=GRAY,
                         align=PP_ALIGN.CENTER)

        # 2×2 cells
        default_colors = [BLUE_LT, GREEN_LT, ORANGE_LT, PURPLE_LT]
        for row in range(2):
            for col_idx in range(2):
                cx = grid_x + col_idx * (cell_w + cell_gap)
                cy = grid_y + row * (cell_h + cell_gap)
                cell = cells[row][col_idx] if row < len(cells) and col_idx < len(cells[row]) else {}

                if isinstance(cell, dict):
                    label = cell.get("label", "")
                    desc = cell.get("desc", "")
                    fill = cell.get("color", default_colors[row * 2 + col_idx])
                elif isinstance(cell, (list, tuple)):
                    label = cell[0] if len(cell) > 0 else ""
                    desc = cell[1] if len(cell) > 1 else ""
                    fill = cell[2] if len(cell) > 2 else default_colors[row * 2 + col_idx]
                else:
                    label, desc, fill = str(cell), "", default_colors[row * 2 + col_idx]

                if isinstance(fill, str):
                    fill = resolve_named_color(fill)

                text = f"{label}\n\n{desc}" if desc else label
                self.box(slide, cx, cy, cell_w, cell_h, fill, None,
                         text, sz=SZ_BODY, tc=_auto_tc(fill),
                         align=PP_ALIGN.CENTER, anchor="ctr")
        return slide

    # ─── Stat (big numbers) ────────────────────────────────

    def stat_slide(self, title, stats, subtitle=None):
        """Big number statistics display.
        stats: [{"value": "99.9%", "label": "Uptime", "color": "BLUE"}, ...]
        Or: [["99.9%", "Uptime", "BLUE"], ...]
        """
        slide = self._new_slide(title)
        sub_offset = 0
        if subtitle:
            self.textbox(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 0.38,
                         subtitle, sz=SZ_BODY, color=GRAY)
            sub_offset = 0.55

        n = len(stats)
        zone_top = CONTENT_TOP + sub_offset
        zone_h = CONTENT_H - sub_offset
        gap = 0.30
        cw = (CONTENT_W - gap * (n - 1)) / n
        ch = min(3.2, zone_h * 0.65)
        cy = zone_top + (zone_h - ch) / 2

        for i, s in enumerate(stats):
            if isinstance(s, dict):
                value = s.get("value", "")
                label = s.get("label", "")
                color = s.get("color", BLUE)
            elif isinstance(s, (list, tuple)):
                value = s[0] if len(s) > 0 else ""
                label = s[1] if len(s) > 1 else ""
                color = s[2] if len(s) > 2 else BLUE

            if isinstance(color, str):
                color = resolve_named_color(color)

            cx = MARGIN_X + i * (cw + gap)

            # Background card
            self.box(slide, cx, cy, cw, ch, WHITE, GRAY_LT, "",
                     anchor="ctr")

            # Accent top bar
            self.rect(slide, cx, cy, cw, 0.06, color)

            # Big value
            value_h = 1.2
            vy = cy + (ch - value_h - 0.6) / 2 + 0.15
            self.textbox(slide, cx + 0.15, vy, cw - 0.30, value_h,
                         value, sz=Pt(44), bold=True, color=color,
                         align=PP_ALIGN.CENTER)

            # Label below
            self.textbox(slide, cx + 0.15, vy + value_h, cw - 0.30, 0.50,
                         label, sz=SZ_BODY, color=GRAY,
                         align=PP_ALIGN.CENTER)
        return slide

    # ─── Two-Column ─────────────────────────────────────────

    def two_column_slide(self, title, left_content, right_content):
        slide = self._new_slide(title)
        for i, content in enumerate([left_content, right_content]):
            x = COL2_LEFT_X if i == 0 else COL2_RIGHT_X
            self._render_content(slide, x, CONTENT_TOP, COL2_W,
                                 CONTENT_H, content)
        return slide

    # ─── Custom ─────────────────────────────────────────────

    def custom_slide(self, title, layout="plain"):
        layouts = {
            "plain": self.L_PLAIN,
            "logo": self.L_TONLY,
            "content": self.L_CONTENT,
            "blank": self.L_BLANK,
        }
        slide = self.prs.slides.add_slide(layouts.get(layout, self.L_PLAIN))
        page = self._next_page()
        if layout != "blank":
            self._set_title(slide, title)
            self._set_footer(slide, page)
        return slide

    # ═══════════════════════════════════════════════════════════
    # SAVE
    # ═══════════════════════════════════════════════════════════

    def save(self, output_path):
        self.prs.save(output_path)
        n = len(self.prs.slides)
        print(f"Saved: {output_path} ({n} slides)")
        return output_path

    # ═══════════════════════════════════════════════════════════
    # PRIVATE
    # ═══════════════════════════════════════════════════════════

    def _render_content(self, slide, x, y, w, h, content):
        ctype = content.get("type", "bullets")
        if ctype == "bullets":
            items = content.get("items", [])
            col_title = content.get("title")
            header = content.get("header")
            hc = content.get("header_color", BLUE)
            cy = y

            if header:
                hh = 0.50
                self.box(slide, x, cy, w, hh, hc, None, header,
                         sz=SZ_BODY, bold=True, tc=WHITE)
                cy += hh + 0.12

            bullet_h = h - (cy - y)
            self.bullet_box(slide, x, cy, w, bullet_h, items,
                            sz=SZ_BODY, title=col_title,
                            title_color=content.get("title_color", DARK))

        elif ctype == "box":
            text = content.get("text", "")
            fill = content.get("fill", GRAY_LT)
            border = content.get("border", None)
            self.box(slide, x, y, w, h, fill, border, text,
                     sz=SZ_BODY, tc=content.get("color", BLACK),
                     align=PP_ALIGN.LEFT)
